"""Shared helpers for building FeeBridge slide decks with python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0F, 0x18, 0x30)
BLUE = RGBColor(0x2B, 0x59, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
GREY = RGBColor(0x6B, 0x76, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def new_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_slide(prs, title, subtitle):
    slide = _blank(prs)
    band = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(46)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(22)
    r2.font.color.rgb = RGBColor(0x9F, 0xB0, 0xE0)

    accent = slide.shapes.add_shape(1, Inches(0.9), Inches(2.25), Inches(2.2), Inches(0.10))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()
    return slide


def content_slide(prs, title, bullets, subtitle=None):
    """bullets: list of (text, level) tuples; level 0 = top bullet."""
    slide = _blank(prs)
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.22), Inches(12), Inches(0.9))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        sp = tf.add_paragraph()
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(13)
        sr.font.color.rgb = RGBColor(0x9F, 0xB0, 0xE0)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.6))
    bf = body.text_frame
    bf.word_wrap = True
    first = True
    for text, level in bullets:
        p = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        run.text = ("• " if level == 0 else "– ") + text
        run.font.size = Pt(18 if level == 0 else 15)
        run.font.color.rgb = NAVY if level == 0 else GREY
        p.space_after = Pt(7)
    return slide


def save(prs, path):
    prs.save(path)
    print("wrote", path)
